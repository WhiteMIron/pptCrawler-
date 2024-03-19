from pptx import Presentation
import openpyxl
import copy
from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
import re 
# Type1.점검일지

fileName = "대검 점검일지 영상물중계_2024"
def extract_data_from_tables(slide, extracted_data):
    found_words = set()  # 이미 발견된 단어를 기록하기 위한 집합
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row in range(len(table.rows)):
               for col in range(len(table.columns)):
                cell = table.cell(row, col)
                text = cell.text.strip()
                if text in extracted_data and text not in found_words:
                        model_cell = table.cell(row, col+1)
                        data =model_cell.text.strip()
                        extracted_data[text] = data
                        found_words.add(text)  # 해당 단어를 발견한 후에 집합에 추가
    return extracted_data


def extract_data_from_presentation(presentation,base_data):
    data_per_slide = []
    for slide in presentation.slides:
        extracted_data = copy.copy(base_data)
        extracted_data = extract_data_from_tables(slide,extracted_data)
        data_per_slide.append(copy.deepcopy(extracted_data))
    return data_per_slide

def write_to_excel(data_per_slide):
    workbook = openpyxl.Workbook()

    sheet = workbook.active

    # 엑셀 헤더 작성
    headers = ["HostName", "Model", "Vendor","S/N","OS","CPU model","DISK","Memory(GB)"]
    sheet.append(headers)
    # 각 슬라이드의 데이터를 엑셀에 쓰기
    for slide_data in data_per_slide:
        row = []
        values_list = list(slide_data.values())  
        for index, header in enumerate(headers):
            processed_value = replace_special_characters(values_list[index]) 
            if(header=="Memory(GB)"):
                 processed_value=mb_to_gb(values_list[index])
            row.append(processed_value)  
        sheet.append(row)

    workbook.save(f"{fileName}.xlsx")
def replace_special_characters(text):
    text = ILLEGAL_CHARACTERS_RE.sub(r'',text)
    return text

def mb_to_gb(mb_string):
    result_search = re.search(r'(\d+)(?i:mb)(?![mb])', mb_string)
    if result_search:
       mb_value = int(result_search.group(1))
       gb_value = mb_value / 1024  
       return f"{gb_value:.0f}GB"
    else:
        result_search = re.search(r'(\d+)(?i:g)\b', mb_string)
        if result_search:
          converted_string = re.sub(r'(\d+)(?i:g)\b', r'\1GB', mb_string)
          return converted_string 
        else :
           return mb_string

def main():
    presentation = Presentation(f"{fileName}.pptx") 
    base_data = {"Hostname": '', "Model": '', "Vendor": '', "S/N": '', "OS": '', "CPU model": '', "DISK": '', "MEM": ''}
    data_per_slide = extract_data_from_presentation(presentation, base_data)
  
    write_to_excel(data_per_slide)

if __name__ == "__main__":
    main()
 