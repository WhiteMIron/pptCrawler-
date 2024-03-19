from pptx import Presentation
import openpyxl
import copy
# Type1.점검일지
filename="1.대검 점검일지_2024"
def extract_data_from_tables(slide, extracted_data):
    found_words = set()  # 이미 발견된 단어를 기록하기 위한 집합
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row in range(len(table.rows)):
               for col in range(len(table.columns)):
                cell = table.cell(row, col)
                text = cell.text.strip()
                if text in extracted_data and text not in found_words:
                       if row + 1 < len(table.rows): 
                            model_cell = table.cell(row + 1, col)
                            data =model_cell.text.strip()
                            extracted_data[text] = data
                            found_words.add(text)  # 해당 단어를 발견한 후에 집합에 추가
    return extracted_data


def extract_data_from_presentation(presentation,base_data):
    
    data_per_slide = []
    for slide in presentation.slides:
        extracted_data = copy.copy(base_data)
        extracted_data = extract_data_from_tables(slide,extracted_data)
        data_per_slide.append(copy.deepcopy(extracted_data))
    return data_per_slide

def write_to_excel(data_per_slide):
    workbook = openpyxl.Workbook()
    sheet = workbook.active

    # 엑셀 헤더 작성
    headers = ["HostName", "Model", "Vendor","S/N","OS","CPU model","DISK","Memory(GB)"]
    sheet.append(headers)
    # 각 슬라이드의 데이터를 엑셀에 쓰기
    for slide_data in data_per_slide:
        row = []
        values_list = list(slide_data.values())  
        for index, header in enumerate(headers):
            row.append(values_list[index])  
        sheet.append(row)

    workbook.save(f"{filename}.xlsx")

def main():
    presentation = Presentation("1.대검 점검일지_2024.pptx") 
    # extracted_data = {"Hostname": '', "장비모델": '', "Vendor": '', "S/N": '', "OS": '', "CPU model": '',"DISK":'', "MEM": ''}
    base_data = {"Hostname": '', "장비모델": '', "Vendor": '', "S/N": '', "OS": '', "CPU model": '', "DISK": '', "MEM": ''}

    data_per_slide = extract_data_from_presentation(presentation, base_data)
  
    write_to_excel(data_per_slide)

if __name__ == "__main__":
    main()
 